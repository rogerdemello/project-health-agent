"""Professional PowerPoint generator for monthly executive presentation."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from loguru import logger

from app.models import WeeklyReport, RiskAssessment, TrendPoint, ProjectTrend
from app.config import OUTPUT_DIR


# Color palette
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF7)
GRAY_TEXT = RGBColor(0x7F, 0x8C, 0x8D)
GREEN = RGBColor(0x27, 0xAE, 0x60)
AMBER = RGBColor(0xF3, 0x9C, 0x12)
RED = RGBColor(0xE7, 0x4C, 0x3C)


def _add_bg(slide, color=LIGHT_GRAY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_bar(slide, title, subtitle=""):
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.font.bold = True

    if subtitle:
        tbox = slide.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(8), Inches(0.35))
        p2 = tbox.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(12)
        p2.font.color.rgb = GRAY_TEXT
        p2.font.italic = True


def _add_box(slide, left, top, width, height):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(0.5)
    return shape


def _add_text(slide, left, top, width, height, text, size=12, bold=False, color=DARK_BLUE, align=PP_ALIGN.LEFT):
    tbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def _add_rag_badge(slide, left, top, rag, size=0.6):
    colors = {"Green": GREEN, "Amber": AMBER, "Red": RED}
    c = colors.get(rag, GRAY_TEXT)
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(size), Inches(size * 0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = c
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = rag
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def _rag_color(score: float) -> RGBColor:
    if score >= 85:
        return GREEN
    elif score >= 60:
        return AMBER
    return RED


def generate_presentation(
    reports: list[WeeklyReport],
    assessments: list[RiskAssessment],
    trends: list[ProjectTrend],
    output_path: str | None = None,
) -> Path:
    """Generate a 6-slide monthly executive presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    if output_path is None:
        output_path = str(OUTPUT_DIR / "Monthly_Executive_Review.pptx")

    # ---- Slide 1: Title ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, DARK_BLUE)
    _add_text(slide, 1, 1.5, 8, 1, "Project Health Review", size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(slide, 1, 2.8, 8, 0.5, "Monthly Executive Summary", size=18, color=RGBColor(0xBD, 0xC3, 0xC7), align=PP_ALIGN.CENTER)
    _add_text(slide, 1, 3.6, 8, 0.4, "Professional Services - AI Agent Generated", size=13, color=RGBColor(0x95, 0xA5, 0xA6), align=PP_ALIGN.CENTER)
    if reports:
        _add_text(slide, 1, 4.3, 8, 0.4, f"Reporting Period: {reports[0].assessment_date}", size=12, color=RGBColor(0x95, 0xA5, 0xA6), align=PP_ALIGN.CENTER)

    # ---- Slide 2: Executive Summary ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_title_bar(slide, "Executive Summary", "Portfolio health overview")

    y = 1.3
    for report in reports:
        _add_box(slide, 0.5, y, 9, 0.6)
        _add_text(slide, 0.7, y + 0.05, 4, 0.3, report.project_name[:50], size=13, bold=True)
        _add_rag_badge(slide, 5.5, y + 0.1, report.overall_status)
        _add_text(slide, 6.3, y + 0.1, 1.2, 0.3, f"{report.risk_score:.0f}/100", size=11, bold=True, color=_rag_color(report.risk_score))
        _add_text(slide, 7.5, y + 0.1, 2, 0.3, f"Score: {report.risk_score:.0f}", size=11, color=GRAY_TEXT)
        y += 0.85

    # Summary insight box
    _add_box(slide, 0.5, y + 0.2, 9, 1)
    green = sum(1 for r in reports if r.overall_status == "Green")
    amber = sum(1 for r in reports if r.overall_status == "Amber")
    red = sum(1 for r in reports if r.overall_status == "Red")
    insight = f"Portfolio: {green} Green, {amber} Amber, {red} Red."
    if red > 0:
        insight += " Immediate attention required on at-risk projects."
    elif amber > 0:
        insight += " Some projects need monitoring and support."
    else:
        insight += " Portfolio is healthy."
    _add_text(slide, 0.8, y + 0.35, 8.5, 0.5, insight, size=12, bold=True, color=DARK_BLUE)

    # ---- Slide 3: Portfolio Health ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_title_bar(slide, "Portfolio Health", "Current scores and trends")

    # Trend chart using latest scores
    chart_data = CategoryChartData()
    chart_data.categories = [r.project_name[:20] for r in reports]
    chart_data.add_series("Risk Score", [r.risk_score for r in reports])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.3), Inches(5.5), Inches(3.5), chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    plot = chart.plots[0]
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = ACCENT_BLUE

    # Trend arrows on right
    y = 1.4
    for trend in trends:
        delta = trend.score_delta
        arrow = "(+)" if delta > 1 else ("(-)" if delta < -1 else "(~)")
        color = GREEN if delta > 1 else (RED if delta < -1 else AMBER)
        _add_text(slide, 6.5, y, 3, 0.3, f"{trend.project_name[:25]}", size=11, bold=True)
        _add_text(slide, 6.5, y + 0.3, 3, 0.3, f"Score: {trend.current_score:.0f}  {arrow} ({delta:+.1f})", size=10, color=color)
        _add_rag_badge(slide, 9, y, trend.current_rag, size=0.5)
        y += 0.7

    # ---- Slide 4: Top Risks & Emerging Themes ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_title_bar(slide, "Top Risks & Emerging Themes", "Cross-project risk identification")

    # Collect all risks
    all_risks: list[str] = []
    for r in reports:
        for risk in r.top_risks:
            label = f"[{r.project_name[:20]}] {risk}"
            if label not in all_risks:
                all_risks.append(label)

    # Emerging themes (cross-project)
    themes = [
        "Schedule pressure is the dominant risk across implementations.",
        "Milestone overruns cluster in the final phases of delivery.",
        "Stakeholder comment volume correlates with project health decline.",
        "Resource contention emerges where projects share specialist pools.",
    ]

    y = 1.3
    _add_text(slide, 0.5, y, 4.5, 0.3, "Top Risks", size=13, bold=True, color=DARK_BLUE)
    _add_text(slide, 5.5, y, 4.5, 0.3, "Emerging Themes", size=13, bold=True, color=DARK_BLUE)

    y += 0.4
    for i, risk in enumerate(all_risks[:6]):
        c = RED if i < 2 else AMBER
        shape = slide.shapes.add_shape(9, Inches(0.7), Inches(y + 0.04), Inches(0.1), Inches(0.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = c
        shape.line.fill.background()
        _add_text(slide, 1.0, y, 4, 0.35, risk[:60], size=10, color=DARK_BLUE)
        if i < len(themes):
            _add_text(slide, 5.7, y, 4, 0.35, themes[i][:70], size=10, color=DARK_BLUE)
        y += 0.4

    # ---- Slide 5: Recommendations ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_title_bar(slide, "Recommendations", "Executive-level action items")

    recs = [
        "Prioritize critical-path recovery for Red-rated projects.",
        "Establish weekly milestone tracking with automated alerts.",
        "Conduct cross-project resource review to resolve specialist contention.",
        "Implement standardized stakeholder check-ins across all projects.",
        "Deploy variance early-warning system (threshold: -5 days).",
    ]

    y = 1.4
    for i, rec in enumerate(recs, 1):
        _add_box(slide, 0.5, y, 9, 0.55)
        _add_text(slide, 0.8, y + 0.08, 0.3, 0.3, str(i), size=14, bold=True, color=ACCENT_BLUE)
        _add_text(slide, 1.3, y + 0.08, 8, 0.4, rec, size=12, color=DARK_BLUE)
        y += 0.65

    # ---- Slide 6: Project Snapshot ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_title_bar(slide, "Project Snapshot", "Detailed view")

    # Table-like layout
    y = 1.3
    headers = ["Project", "Status", "Score", "Schedule", "Milestones", "Critical", "Variance"]
    col_widths = [2.5, 0.8, 0.8, 1.2, 1.2, 1.0, 1.0]
    x = 0.5

    for h, w in zip(headers, col_widths):
        _add_text(slide, x, y, w, 0.3, h, size=10, bold=True, color=WHITE)
        x += w
    # header bg
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(y - 0.05), Inches(8.5), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    # Re-draw header text on top
    x = 0.5
    for h, w in zip(headers, col_widths):
        _add_text(slide, x, y - 0.02, w, 0.3, h, size=10, bold=True, color=WHITE)
        x += w

    y += 0.4
    for report, assessment in zip(reports, assessments):
        dims = {d.name: d.score for d in assessment.dimensions}
        vals = [report.project_name[:30], report.overall_status, f"{report.risk_score:.0f}",
                f"{dims.get('Schedule Health', 0):.0f}", f"{dims.get('Milestones', 0):.0f}",
                f"{dims.get('Critical Tasks', 0):.0f}", f"{dims.get('Variance', 0):.0f}"]
        x = 0.5
        for v, w in zip(vals, col_widths):
            _add_text(slide, x, y, w, 0.3, v, size=9, color=DARK_BLUE)
            x += w
        y += 0.3

    _add_text(slide, 0.5, y + 0.3, 9, 0.3,
              "Auto-generated by Project Health AI Agent", size=9, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

    prs.save(output_path)
    logger.info(f"Presentation saved: {output_path}")
    return Path(output_path)
